"""Document parsers that normalize files into page text records."""

from __future__ import annotations

from pathlib import Path


Page = dict[str, str | int]


def parse_document(path: str | Path) -> list[Page]:
    """Parse a supported document into ``{"text": str, "page": int}`` records."""

    file_path = Path(path)
    suffix = file_path.suffix.lower()

    if suffix == ".pdf":
        return parse_pdf(file_path)
    if suffix == ".docx":
        return parse_docx(file_path)
    if suffix == ".pptx":
        return parse_pptx(file_path)
    if suffix in {".txt", ".md"}:
        return parse_text_file(file_path)

    raise ValueError(f"Unsupported document type: {file_path.suffix}")


def parse_pdf(path: str | Path) -> list[Page]:
    """Parse PDF pages with PyMuPDF, falling back to pypdf."""

    file_path = Path(path)
    try:
        return _parse_pdf_pymupdf(file_path)
    except Exception as pymupdf_error:
        try:
            return _parse_pdf_pypdf(file_path)
        except Exception as pypdf_error:
            raise RuntimeError(
                f"Could not parse PDF {file_path} with PyMuPDF or pypdf."
            ) from pypdf_error


def _parse_pdf_pymupdf(path: Path) -> list[Page]:
    import fitz

    pages: list[Page] = []
    with fitz.open(path) as document:
        for index, page in enumerate(document, start=1):
            text = page.get_text() or ""
            if text.strip():
                pages.append({"text": text, "page": index})
    return pages


def _parse_pdf_pypdf(path: Path) -> list[Page]:
    from pypdf import PdfReader

    pages: list[Page] = []
    reader = PdfReader(str(path))
    for index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            pages.append({"text": text, "page": index})
    return pages


def parse_docx(path: str | Path) -> list[Page]:
    """Parse a DOCX document into a single page-like text record."""

    from docx import Document

    document = Document(str(path))
    text = "\n".join(paragraph.text for paragraph in document.paragraphs)
    return _single_page(text)


def parse_pptx(path: str | Path) -> list[Page]:
    """Parse a PPTX document into one text record per slide."""

    from pptx import Presentation

    presentation = Presentation(str(path))
    pages: list[Page] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
        text = "\n".join(parts)
        if text.strip():
            pages.append({"text": text, "page": index})
    return pages


def parse_text_file(path: str | Path) -> list[Page]:
    """Parse a TXT or Markdown file into a single page-like text record."""

    text = Path(path).read_text(encoding="utf-8")
    return _single_page(text)


def _single_page(text: str) -> list[Page]:
    if not text.strip():
        return []
    return [{"text": text, "page": 1}]
